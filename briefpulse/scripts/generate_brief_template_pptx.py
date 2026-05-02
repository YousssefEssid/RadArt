"""
Génère le modèle de brief : une diapo par rubrique, sans exemple de contenu,
seulement le titre de section et une consigne indiquant quoi saisir.
Usage : python scripts/generate_brief_template_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "frontend" / "public" / "templates" / "brief-modele-a-remplir.pptx"

# (titre de la diapo, consigne pour le champ — pas de contenu d'exemple)
SECTIONS: list[tuple[str, str]] = [
    (
        "Client et marque",
        "À remplir : nom du client, nom de la marque, périmètre (pays, BU…).",
    ),
    (
        "Contexte et objectifs",
        "À remplir : situation actuelle, enjeu business, objectifs mesurables de la campagne.",
    ),
    (
        "Cible et insights",
        "À remplir : qui viser (segments), comportements, insights ou vérités consommateur.",
    ),
    (
        "Message et territoire créatif",
        "À remplir : promesse ou message clé, tonalité, univers créatif souhaité (humour, émotion…).",
    ),
    (
        "Canaux et livrables",
        "À remplir : médias prioritaires (TV, digital, social…), formats et livrables attendus.",
    ),
    (
        "Budget et calendrier",
        "À remplir : budget indicatif ou fourchette, dates clés, deadline de réponse.",
    ),
    (
        "Contraintes",
        "À remplir : obligations légales, sujets à éviter, concurrents, autres contraintes.",
    ),
]


def main() -> None:
    prs = Presentation()

    # Page d'accueil
    layout_title = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(layout_title)
    s0.shapes.title.text = "Brief créatif — modèle à compléter"
    subtitle = s0.placeholders[1]
    subtitle.text = (
        "Chaque diapo indique uniquement quoi saisir dans cette partie du brief. "
        "Remplacez la consigne par votre texte, ou rédigez en dessous dans la même zone."
    )

    layout_section = prs.slide_layouts[1]
    for heading, instruction in SECTIONS:
        slide = prs.slides.add_slide(layout_section)
        slide.shapes.title.text = heading
        body = slide.placeholders[1].text_frame
        body.text = instruction
        for p in body.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Written {OUT}")


if __name__ == "__main__":
    main()
