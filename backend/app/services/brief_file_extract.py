"""Extraction de texte depuis un brief importé (PowerPoint, Word, PDF, texte)."""
from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path


def extract_text_from_filename_and_bytes(filename: str, content: bytes) -> str:
    if not content:
        raise ValueError("Fichier vide.")
    ext = Path(filename or "").suffix.lower()
    if ext in (".txt", ".md", ".csv"):
        t = content.decode("utf-8", errors="replace").strip()
        if not t:
            raise ValueError("Fichier texte vide.")
        return t

    if ext == ".pdf":
        return _extract_pdf(content)
    if ext == ".docx":
        return _extract_docx(content)
    if ext in (".pptx",):
        return _extract_pptx(content)
    if ext == ".ppt":
        # Ancien binaire OLE ; python-pptx ne lit pas ce format.
        raise ValueError(
            "Le format .ppt (PowerPoint 97-2003) ne permet pas l'extraction automatique ici. "
            "Enregistrez le fichier au format .pptx (Fichier, Enregistrer sous, type PowerPoint .pptx), "
            "puis importez a nouveau, ou copiez-collez le texte des diapositives dans le champ Brief."
        )

    raise ValueError(
        f"Extension non prise en charge ({ext or 'inconnue'}). Utilisez .pptx, .docx, .pdf ou .txt."
    )


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(content))
    chunks: list[str] = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            chunks.append(t.strip())
    text = "\n\n".join(chunks).strip()
    if not text:
        raise ValueError("PDF sans texte extractible (image scannée ?). Copiez le texte manuellement.")
    return text


def _extract_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(BytesIO(content))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paras)
    if not text:
        raise ValueError("Document Word sans paragraphes texte détectés.")
    return text


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    try:
        prs = Presentation(BytesIO(content))
    except PackageNotFoundError as e:
        raise ValueError(
            "Fichier PowerPoint illisible ou ancien format .ppt. Enregistrez au format .pptx puis reessayez."
        ) from e

    texts: list[str] = []

    def collect_from_shapes(shapes) -> None:
        for shape in shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
            if getattr(shape, "shape_type", None) == 6:  # GROUP
                try:
                    collect_from_shapes(shape.shapes)
                except Exception:
                    pass

    for slide in prs.slides:
        collect_from_shapes(slide.shapes)

    # Dédupliquer lignes répétées (masters)
    lines_out: list[str] = []
    seen: set[str] = set()
    for line in texts:
        key = re.sub(r"\s+", " ", line.lower())
        if key in seen:
            continue
        seen.add(key)
        lines_out.append(line)
    text = "\n".join(lines_out).strip()
    if not text:
        raise ValueError("Présentation sans texte extractible. Copiez le contenu des diapositives dans la zone Brief.")
    return text
